#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
build_deck.py — outline.json → .pptx 渲染器（python-pptx，无 AI API 依赖）

用法：
    python build_deck.py outline.json -o presentation.pptx
    python build_deck.py outline.json -o out.pptx --preset cns-bio-light

吸收 siril9/presentation-skill 的 source-first 范式 + davila7 的 DESIGN 常量。
"""
import argparse, json, sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- Preset（生信专用 cns-bio-light）----
# font_name 必须设——python-pptx 不设则用默认主题字体（Calibri），不含中文字形 → 豆腐块。
# 'Microsoft YaHei' 在 Windows PowerPoint 里中文正常；macOS 会回退到 PingFang。
# 英文期刊投稿最终版：改成 'Arial' 并把中文翻译为英文。
PRESETS = {
    "cns-bio-light": {
        "bg": RGBColor(0xFF,0xFF,0xFF),
        "title": RGBColor(0x1F,0x3A,0x5F),    # Navy
        "accent": RGBColor(0x3D,0x7A,0xAB),    # Blue
        "body": RGBColor(0x33,0x33,0x33),       # Dark grey
        "caption": RGBColor(0x66,0x66,0x66),    # Light grey
        "pass": RGBColor(0x2E,0x8B,0x57),       # Green
        "fail": RGBColor(0xE2,0x5D,0x5D),       # Red
        "warn": RGBColor(0xE8,0xA8,0x38),       # Orange
        "font_name": "Microsoft YaHei",         # CJK-safe; macOS auto-fallback to PingFang
    }
}

def build(outline_path, output_path, preset_name="cns-bio-light"):
    outline = json.loads(Path(outline_path).read_text(encoding="utf-8"))
    preset = PRESETS.get(preset_name, PRESETS["cns-bio-light"])
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for sdef in outline.get("slides", []):
        variant = sdef.get("variant","bullets")
        title = sdef.get("title","")
        s = prs.slides.add_slide(blank)
        bg = s.background.fill; bg.solid(); bg.fore_color.rgb = preset["bg"]
        if variant == "title":
            _title_slide(s, sdef, preset)
        elif variant == "section":
            _section_slide(s, title, preset)
        elif variant == "figure-hero":
            _figure_hero(s, sdef, preset)
        elif variant == "figure-sidebar":
            _figure_sidebar(s, sdef, preset)
        elif variant == "figure-dual":
            _figure_dual(s, sdef, preset)
        elif variant == "figure-top-text":
            _figure_top_text(s, sdef, preset)
        elif variant == "figure-grid":
            _figure_grid(s, sdef, preset)
        elif variant == "split-compare":
            _split_compare(s, sdef, preset)
        elif variant in ("scientific-figure", "image-sidebar"):
            _figure_sidebar(s, sdef, preset)  # 兼容旧名
        elif variant == "results-table":
            _table_slide(s, sdef, preset)
        elif variant == "methods-flow":
            _flow_slide(s, sdef, preset)
        else:
            _bullets_slide(s, sdef, preset)

    prs.save(str(output_path))
    print(f"SAVED {output_path} ({len(prs.slides)} slides)")

def _add_text(s, left, top, width, height, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT, font_name=None):
    tb = s.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.color.rgb = color
    p.font.bold = bold; p.font.italic = italic; p.alignment = align
    if font_name: p.font.name = font_name   # CJK-safe
    return tb

def _add_bullets(s, left, top, width, height, bullets, size, color, font_name=None):
    tb = s.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text = "• " + str(b); p.font.size = Pt(size); p.font.color.rgb = color
        if font_name: p.font.name = font_name
        p.space_after = Pt(8)

# ---- 安全区域（防重叠）----
SLIDE_W, SLIDE_H = 13.333, 7.5
SAFE_GAP = 0.3  # inch, 图片与文字之间的最小安全间距
TITLE_TOP, TITLE_H = 0.3, 0.7   # 标题区
CONTENT_TOP = 1.2                # 内容区起点
CAPTION_TOP = 6.7                # 图注区起点
CONTENT_H = CAPTION_TOP - CONTENT_TOP - 0.2  # 可用内容高度

def _add_title(s, title, preset):
    """统一标题放置：顶部 0.3-1.0inch，绝不与内容区重叠"""
    if title:
        _add_text(s, Inches(0.5), Inches(TITLE_TOP), Inches(SLIDE_W-1), Inches(TITLE_H),
                  title, 28, preset["title"], bold=True, font_name=preset["font_name"])

def _add_caption(s, caption, preset, left=0.5, width=None):
    """统一图注放置：底部 6.7-7.2inch"""
    if caption:
        w = width or (SLIDE_W - 1)
        _add_text(s, Inches(left), Inches(CAPTION_TOP), Inches(w), Inches(0.5),
                  caption, 10, preset["caption"], italic=True, font_name=preset["font_name"])

def _place_image(s, img_path, left, top, max_w, max_h):
    """放置图片并自动缩放到 max_w × max_h 内（保持宽高比）"""
    pic = s.shapes.add_picture(img_path, Inches(left), Inches(top))
    w_in = pic.width / 914400  # EMU → inch
    h_in = pic.height / 914400
    # 缩放到 max_w × max_h
    ratio = min(max_w / w_in, max_h / h_in, 1.0)  # 不放大只缩小
    pic.width = int(pic.width * ratio)
    pic.height = int(pic.height * ratio)
    return pic

def _title_slide(s, d, preset):
    _add_text(s, Inches(1), Inches(2.5), Inches(11), Inches(1.5),
              d.get("title",""), 40, preset["title"], bold=True, align=PP_ALIGN.CENTER, font_name=preset["font_name"])
    if d.get("subtitle"):
        _add_text(s, Inches(1), Inches(4), Inches(11), Inches(1),
                  d["subtitle"], 20, preset["caption"], align=PP_ALIGN.CENTER, font_name=preset["font_name"])

def _section_slide(s, title, preset):
    _add_text(s, Inches(1), Inches(3), Inches(11), Inches(1.5),
              title, 36, preset["title"], bold=True, align=PP_ALIGN.CENTER, font_name=preset["font_name"])

# ---- 新布局：figure-hero（全宽大图）----
def _figure_hero(s, d, preset):
    _add_title(s, d.get("title",""), preset)
    img = d.get("image")
    if img and Path(img).exists():
        # 全宽：0.5 → SLIDE_W-0.5，但限制最大高度 = CONTENT_H
        max_w = SLIDE_W - 1.0
        _place_image(s, img, left=0.5, top=CONTENT_TOP, max_w=max_w, max_h=CONTENT_H)
    _add_caption(s, d.get("caption"), preset)
    # hero 模式：文字放到图注下方（如果有的话）——不与图重叠
    # 如果有 bullets，建议用下一张 slide 而不是塞进来

# ---- 新布局：figure-dual（左右双图对比）----
def _figure_dual(s, d, preset):
    _add_title(s, d.get("title",""), preset)
    img1 = d.get("image"); img2 = d.get("image2")
    half_w = (SLIDE_W - 1.0 - SAFE_GAP) / 2  # 每半宽度
    if img1 and Path(img1).exists():
        _place_image(s, img1, left=0.5, top=CONTENT_TOP, max_w=half_w, max_h=CONTENT_H)
        _add_text(s, Inches(0.5), Inches(CAPTION_TOP - 0.4), Inches(half_w), Inches(0.35),
                  d.get("caption_left",""), 9, preset["caption"], italic=True, font_name=preset["font_name"])
    if img2 and Path(img2).exists():
        left2 = 0.5 + half_w + SAFE_GAP
        _place_image(s, img2, left=left2, top=CONTENT_TOP, max_w=half_w, max_h=CONTENT_H)
        _add_text(s, Inches(left2), Inches(CAPTION_TOP - 0.4), Inches(half_w), Inches(0.35),
                  d.get("caption_right",""), 9, preset["caption"], italic=True, font_name=preset["font_name"])
    _add_caption(s, d.get("caption"), preset)

# ---- 新布局：figure-top-text（图上文下）----
def _figure_top_text(s, d, preset):
    _add_title(s, d.get("title",""), preset)
    # 文字区：CONTENT_TOP → CONTENT_TOP + 2.2inch
    text_h = 2.2
    if d.get("bullets"):
        _add_bullets(s, Inches(0.8), Inches(CONTENT_TOP), Inches(SLIDE_W-1.6), Inches(text_h),
                     d["bullets"], 16, preset["body"], font_name=preset["font_name"])
    # 图区：文字底部 + SAFE_GAP → CAPTION_TOP - 0.2
    img_top = CONTENT_TOP + text_h + SAFE_GAP
    img_h = CAPTION_TOP - 0.2 - img_top
    img = d.get("image")
    if img and Path(img).exists():
        _place_image(s, img, left=1.0, top=img_top, max_w=SLIDE_W-2.0, max_h=img_h)
    _add_caption(s, d.get("caption"), preset)

# ---- 新布局：figure-grid（2×2 四宫格）----
def _figure_grid(s, d, preset):
    _add_title(s, d.get("title",""), preset)
    images = d.get("images", [])
    if not images and d.get("image"):
        images = [d["image"]]
    cell_w = (SLIDE_W - 1.0 - SAFE_GAP) / 2
    cell_h = (CONTENT_H - SAFE_GAP) / 2
    positions = [
        (0.5, CONTENT_TOP),
        (0.5 + cell_w + SAFE_GAP, CONTENT_TOP),
        (0.5, CONTENT_TOP + cell_h + SAFE_GAP),
        (0.5 + cell_w + SAFE_GAP, CONTENT_TOP + cell_h + SAFE_GAP),
    ]
    for i, img in enumerate(images[:4]):
        if Path(img).exists():
            left, top = positions[i]
            _place_image(s, img, left=left, top=top, max_w=cell_w, max_h=cell_h)
    _add_caption(s, d.get("caption"), preset)

# ---- 更新后的 figure-sidebar（防重叠版）----
def _figure_sidebar(s, d, preset):
    _add_title(s, d.get("title",""), preset)
    img = d.get("image")
    img_w = SLIDE_W * 0.55 - 0.5  # 图片占 55%，留安全间距
    if img and Path(img).exists():
        _place_image(s, img, left=0.5, top=CONTENT_TOP, max_w=img_w, max_h=CONTENT_H)
    _add_caption(s, d.get("caption"), preset, left=0.5, width=img_w)
    # 文字区：图片右侧 + SAFE_GAP，绝不重叠
    text_left = 0.5 + img_w + SAFE_GAP
    text_w = SLIDE_W - text_left - 0.5
    if d.get("bullets"):
        _add_bullets(s, Inches(text_left), Inches(CONTENT_TOP + 0.3), Inches(text_w), Inches(CONTENT_H - 0.5),
                     d["bullets"], 14, preset["body"], font_name=preset["font_name"])

# ---- 兼容旧版 scientific-figure（路由到 sidebar）----
def _figure_slide(s, d, preset, max_panels=4):
    _figure_sidebar(s, d, preset)  # 统一走 sidebar 布局

# ---- 新布局：split-compare（左右分屏对比）----
def _split_compare(s, d, preset):
    _add_title(s, d.get("title",""), preset)
    half_w = (SLIDE_W - 1.0 - SAFE_GAP) / 2
    # 左半
    left_d = d.get("left", {})
    if left_d.get("title"):
        _add_text(s, Inches(0.5), Inches(CONTENT_TOP), Inches(half_w), Inches(0.4),
                  left_d["title"], 16, preset["accent"], bold=True, font_name=preset["font_name"])
    if left_d.get("image") and Path(left_d["image"]).exists():
        _place_image(s, left_d["image"], left=0.5, top=CONTENT_TOP+0.5, max_w=half_w, max_h=CONTENT_H-1.5)
    if left_d.get("bullets"):
        _add_bullets(s, Inches(0.5), Inches(CAPTION_TOP-1.2), Inches(half_w), Inches(1.0),
                     left_d["bullets"], 11, preset["body"], font_name=preset["font_name"])
    # 右半
    right_d = d.get("right", {})
    right_left = 0.5 + half_w + SAFE_GAP
    if right_d.get("title"):
        _add_text(s, Inches(right_left), Inches(CONTENT_TOP), Inches(half_w), Inches(0.4),
                  right_d["title"], 16, preset["accent"], bold=True, font_name=preset["font_name"])
    if right_d.get("image") and Path(right_d["image"]).exists():
        _place_image(s, right_d["image"], left=right_left, top=CONTENT_TOP+0.5, max_w=half_w, max_h=CONTENT_H-1.5)
    if right_d.get("bullets"):
        _add_bullets(s, Inches(right_left), Inches(CAPTION_TOP-1.2), Inches(half_w), Inches(1.0),
                     right_d["bullets"], 11, preset["body"], font_name=preset["font_name"])

def _table_slide(s, d, preset):
    _add_text(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
              d.get("title",""), 28, preset["title"], bold=True, font_name=preset["font_name"])
    tdef = d.get("table", {})
    headers = tdef.get("headers", []); rows = tdef.get("rows", [])
    if not headers: return
    nrows = len(rows)+1; ncols = len(headers)
    tbl_shape = s.shapes.add_table(nrows, ncols, Inches(1), Inches(1.3), Inches(11), Inches(0.4*nrows))
    tbl = tbl_shape.table
    for j,h in enumerate(headers):
        cell = tbl.cell(0,j); cell.text = str(h)
        cell.text_frame.paragraphs[0].font.size = Pt(12); cell.text_frame.paragraphs[0].font.bold = True; cell.text_frame.paragraphs[0].font.name = preset["font_name"]
        cell.text_frame.paragraphs[0].font.color.rgb = preset["title"]
    for i,row in enumerate(rows):
        for j,val in enumerate(row):
            cell = tbl.cell(i+1,j); cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(11); cell.text_frame.paragraphs[0].font.name = preset["font_name"]
            cell.text_frame.paragraphs[0].font.color.rgb = preset["body"]

def _flow_slide(s, d, preset):
    _add_text(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
              d.get("title",""), 28, preset["title"], bold=True, font_name=preset["font_name"])
    steps = d.get("steps",[])
    if not steps: return
    n = len(steps); step_w = 11.0/n; arrow = " → "
    tb = s.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11), Inches(1))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = arrow.join(steps); p.font.size = Pt(18); p.font.name = preset["font_name"]; p.font.color.rgb = preset["accent"]
    p.alignment = PP_ALIGN.CENTER; p.font.bold = True

def _bullets_slide(s, d, preset):
    _add_text(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
              d.get("title",""), 28, preset["title"], bold=True, font_name=preset["font_name"])
    _add_bullets(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5),
                 d.get("bullets",[]), 18, preset["body"], font_name=preset["font_name"])

def main():
    ap = argparse.ArgumentParser(description="outline.json → .pptx (python-pptx, no AI API)")
    ap.add_argument("outline", help="outline.json path")
    ap.add_argument("-o","--output", default="presentation.pptx", help="output .pptx")
    ap.add_argument("--preset", default="cns-bio-light", choices=list(PRESETS))
    a = ap.parse_args()
    build(a.outline, a.output, a.preset)

if __name__ == "__main__":
    main()
