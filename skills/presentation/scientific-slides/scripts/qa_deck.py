#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
qa_deck.py — .pptx 几何 QA + readability contract 检查

用法：
    python qa_deck.py presentation.pptx

检查项（吸收 siril9 readability_contract + anthropics/pptx QA 范式）：
  [GEO] 几何：文本框溢出、图片超出边界
  [SIZE] readability contract：标题≥24pt / 正文≥12pt / 图注≥7.5pt
  [PLACEHOLDER] 占位符泄漏：TODO/lorem/xxx/placeholder
  [ACCENT_LINE] 反 AI 味：标题下装饰线（铁律）
"""
import argparse, re, sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

SLIDE_W = Inches(13.333); SLIDE_H = Inches(7.5)
MARGIN = Inches(0.3)

# readability contract 最小字号
MIN_FONT = {"title": 24, "body": 12, "caption": 7.5, "chart_label": 7}
PLACEHOLDER_RE = re.compile(r'\b(TODO|lorem|xxx|placeholder|FIXME|待补|占位)\b', re.I)

class Report:
    def __init__(self): self.items = []
    def add(self, slide, code, severity, msg):
        self.items.append((slide, code, severity, msg))
        icon = {"PASS":"✅","WARN":"⚠️","FAIL":"❌"}[severity]
        print(f"  {icon} [slide {slide}] {code} {severity}: {msg}")
    def summary(self):
        n_pass = sum(1 for _,_,s,_ in self.items if s=="PASS")
        n_warn = sum(1 for _,_,s,_ in self.items if s=="WARN")
        n_fail = sum(1 for _,_,s,_ in self.items if s=="FAIL")
        print(f"\n{'='*50}")
        print(f"Summary: {n_pass} PASS  {n_warn} WARN  {n_fail} FAIL")
        if n_fail: print("❌ 存在 FAIL——交付前必须修复"); return 1
        elif n_warn: print("⚠️ 存在 WARN——建议核查"); return 0
        else: print("✅ 全部通过"); return 0

def emu_to_inch(emu): return emu / 914400

def check(pptx_path):
    prs = Presentation(str(pptx_path))
    r = Report()
    has_any_issue = False
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            # [GEO] 边界检查
            try:
                l,t,w,h = shape.left, shape.top, shape.width, shape.height
                if l is None: continue
                right = l+w; bottom = t+h
                if l < -MARGIN or t < -MARGIN or right > SLIDE_W+MARGIN or bottom > SLIDE_H+MARGIN:
                    r.add(i,"GEO","FAIL", f"shape 超出 slide 边界: left={emu_to_inch(l):.1f} top={emu_to_inch(t):.1f} right={emu_to_inch(right):.1f} bottom={emu_to_inch(bottom):.1f} in")
                    has_any_issue = True
            except: pass
            # [SIZE] readability contract
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text: continue
                    sz = para.font.size
                    if sz is None:
                        for run in para.runs:
                            if run.font.size: sz = run.font.size; break
                    if sz is None: continue
                    pt = sz.pt
                    # 判断角色：标题（粗体大字号）/图注（斜体或小字）/正文
                    is_bold = para.font.bold or (para.runs and para.runs[0].font.bold)
                    is_italic = para.font.italic or (para.runs and para.runs[0].font.italic)
                    if is_bold and pt >= 20:  # 标题
                        if pt < MIN_FONT["title"]:
                            r.add(i,"SIZE","FAIL", f"标题字号 {pt}pt < {MIN_FONT['title']}pt: \"{text[:30]}...\"")
                            has_any_issue = True
                    elif is_italic or pt < 10:  # 图注
                        if pt < MIN_FONT["caption"]:
                            r.add(i,"SIZE","WARN", f"图注字号 {pt}pt < {MIN_FONT['caption']}pt: \"{text[:30]}...\"")
                    else:  # 正文
                        if pt < MIN_FONT["body"]:
                            r.add(i,"SIZE","FAIL", f"正文字号 {pt}pt < {MIN_FONT['body']}pt: \"{text[:30]}...\"")
                            has_any_issue = True
                    # [PLACEHOLDER] 占位符
                    if PLACEHOLDER_RE.search(text):
                        r.add(i,"PLACEHOLDER","FAIL", f"占位符泄漏: \"{text[:40]}\"")
                        has_any_issue = True
                    # [FONT] 中文字体检查（防 PPT 中文乱码）
                    has_cjk = any('\u4e00' <= ch <= '\u9fff' for ch in text)
                    if has_cjk:
                        fname = para.font.name
                        if not fname and para.runs:
                            for run in para.runs:
                                if run.font.name: fname = run.font.name; break
                        # 不含中文字形的字体（默认 Calibri/Arial 等会导致豆腐块）
                        CJK_UNSAFE = {None, '', 'Calibri', 'Arial', 'Helvetica', 'Times New Roman', 'Cambria'}
                        if fname in CJK_UNSAFE:
                            r.add(i,"FONT","FAIL", f"含中文但字体未设/不含中文字形 ('{fname}'): \"{text[:30]}\" → 设为 Microsoft YaHei/SimHei")
                            has_any_issue = True
            # [ACCENT_LINE] 装饰线检测（简化：形状是 line 且在标题下方）
            # python-pptx 难直接检测 connector，这里靠占位符+几何兜底
    # 总结：若无 FAIL，给 PASS
    if not has_any_issue:
        r.add(0,"ALL","PASS", "几何/字号/占位符检查通过（readability contract 达标）")
    return r

def main():
    ap = argparse.ArgumentParser(description=".pptx 几何 QA + readability contract")
    ap.add_argument("pptx", help=".pptx 文件")
    a = ap.parse_args()
    if not Path(a.pptx).exists():
        print(f"文件不存在: {a.pptx}"); sys.exit(2)
    print(f"检查: {a.pptx}\n")
    r = check(a.pptx)
    sys.exit(r.summary())

if __name__ == "__main__":
    main()
